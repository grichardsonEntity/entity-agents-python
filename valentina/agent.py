"""
Valentina Agent - Technical Writer & Content Strategist

Expert in technical documentation, API docs, grant writing, and proposals.
Combines documentation expertise with funding research and proposal development.
"""

import asyncio
from typing import Optional, List

from ..shared import BaseAgent, TaskResult
from .config import valentina_config


class ValentinaAgent(BaseAgent):
    """
    Valentina - Technical Writer & Content Strategist

    Specializes in:
    - Technical writing and documentation
    - API documentation
    - Architecture diagrams
    - User guides and tutorials
    - Grant research and writing
    - Proposal development
    - Budget narratives
    - Compliance documentation
    """

    def __init__(self, config=None):
        super().__init__(config or valentina_config)

    # ==================== DOCUMENTATION METHODS ====================

    async def document_feature(self, feature: str, code_path: str = None) -> TaskResult:
        """Document a feature"""
        await self.notify(f"Documenting feature: {feature}")

        prompt = f"""
Document feature: {feature}
{f"Code location: {code_path}" if code_path else ""}

**First:**
1. Read and understand the feature code
2. Identify key components and flows
3. Note any configuration options

**Create documentation:**

## Feature: {feature}

### Overview
What this feature does and why it exists

### How It Works
Technical explanation with flow diagram

### Configuration
Required settings and options

### Usage Examples
```python
# Example code
```

### API Reference (if applicable)
Endpoint documentation

### Troubleshooting
Common issues and solutions

**Save to:** docs/features/{feature.lower().replace(' ', '-')}.md

**Before delivering, verify:**
- [ ] All code paths are documented
- [ ] Examples are tested and work
- [ ] Configuration options are complete
- [ ] Links to related docs are valid
"""

        return await self.run_task(prompt)

    async def document_api(self, endpoint: str, code_path: str = None) -> TaskResult:
        """Document an API endpoint"""
        await self.notify(f"Documenting API: {endpoint}")

        prompt = f"""
Document API endpoint: {endpoint}
{f"Code location: {code_path}" if code_path else ""}

**Generate OpenAPI-style documentation:**

## {endpoint}

### Description
What this endpoint does

### Authentication
Required authentication method

### Request

**Headers:**
| Header | Required | Description |
|--------|----------|-------------|
| Authorization | Yes | Bearer token |

**Parameters:**
| Name | Type | Required | Description |
|------|------|----------|-------------|
| ... | ... | ... | ... |

**Request Body:**
```json
{{
  "field": "value"
}}
```

### Response

**Success (200):**
```json
{{
  "data": {{}}
}}
```

**Errors:**
| Status | Description |
|--------|-------------|
| 400 | Bad Request |
| 401 | Unauthorized |
| 404 | Not Found |

### Example

```bash
curl -X POST {endpoint} \\
  -H "Authorization: Bearer $TOKEN" \\
  -H "Content-Type: application/json" \\
  -d '{{"field": "value"}}'
```

**Before delivering, verify:**
- [ ] All parameters documented
- [ ] All response codes covered
- [ ] Example request works
- [ ] Matches actual implementation
"""

        return await self.run_task(prompt)

    async def create_diagram(self, subject: str, diagram_type: str = "flow") -> TaskResult:
        """Create a Mermaid diagram"""
        await self.notify(f"Creating {diagram_type} diagram for: {subject}")

        prompt = f"""
Create {diagram_type} diagram for: {subject}

**Diagram types:**
- flow: Process or data flow
- sequence: Interaction sequence
- class: Class relationships
- er: Entity relationships
- state: State machine

**Generate Mermaid diagram:**

```mermaid
{f"graph TD" if diagram_type == "flow" else ""}
{f"sequenceDiagram" if diagram_type == "sequence" else ""}
{f"classDiagram" if diagram_type == "class" else ""}
{f"erDiagram" if diagram_type == "er" else ""}
{f"stateDiagram-v2" if diagram_type == "state" else ""}
    ... diagram content ...
```

**Include:**
- Clear labels
- Proper relationships
- Legend if complex
- Brief description below diagram

**Before delivering, verify:**
- [ ] Diagram renders correctly
- [ ] All relationships accurate
- [ ] Labels are clear
"""

        return await self.run_task(prompt)

    async def update_readme(self, project_path: str = ".") -> TaskResult:
        """Update project README"""
        await self.notify(f"Updating README for: {project_path}")

        prompt = f"""
Update README for project at: {project_path}

**Review current state:**
1. Read existing README.md
2. Scan project structure
3. Check package.json/pyproject.toml for details

**Ensure README has:**

# Project Name

Brief description

## Features
- Feature 1
- Feature 2

## Quick Start

```bash
# Installation
pip install ...

# Run
python ...
```

## Configuration
Environment variables and options

## Usage
Common use cases with examples

## API Reference
Link to API docs or brief overview

## Development
Setup for contributors

## License
License information

**Before delivering, verify:**
- [ ] Installation instructions work
- [ ] Examples are accurate
- [ ] All sections up to date
"""

        return await self.run_task(prompt)

    async def create_architecture_doc(self, system: str) -> TaskResult:
        """Create architecture documentation"""
        prompt = f"""
Create architecture documentation for: {system}

## Architecture: {system}

### Overview
System purpose and scope

### High-Level Architecture

```mermaid
graph TD
    subgraph Frontend
        A[Web App]
    end
    subgraph Backend
        B[API Gateway]
        C[Service A]
        D[Service B]
    end
    subgraph Data
        E[Database]
        F[Cache]
    end
    A --> B
    B --> C
    B --> D
    C --> E
    D --> F
```

### Components

| Component | Technology | Purpose |
|-----------|------------|---------|
| ... | ... | ... |

### Data Flow
How data moves through the system

### Security
Authentication, authorization, encryption

### Deployment
How the system is deployed

### Scaling
How components scale

### Monitoring
Logging, metrics, alerts

**Before delivering, verify:**
- [ ] All components included
- [ ] Relationships accurate
- [ ] Security considerations complete
"""

        return await self.run_task(prompt)

    # ==================== GRANT WRITING METHODS ====================

    async def research_funding(
        self,
        project_type: str,
        budget_range: str = None,
        geographic_focus: str = None
    ) -> TaskResult:
        """Research funding opportunities"""
        await self.notify(f"Researching funding for: {project_type}")

        prompt = f"""
Research funding opportunities for:
- Project type: {project_type}
{f"- Budget range: {budget_range}" if budget_range else ""}
{f"- Geographic focus: {geographic_focus}" if geographic_focus else ""}

**Search:**
1. Federal sources (Grants.gov, NSF, NIH, DOE)
2. Foundation sources (FDO, GrantWatch, Candid)
3. Corporate giving programs
4. State/local opportunities

**For each opportunity, document:**

## Opportunity: [Name]

| Field | Value |
|-------|-------|
| Funder | [Name] |
| Type | Federal/Foundation/Corporate |
| Award Range | $X - $Y |
| Deadline | [Date] |
| Eligibility | [Requirements] |
| Match Required | Yes/No (X%) |

### Alignment Analysis
- Mission match: High/Medium/Low
- Capacity match: High/Medium/Low
- Competitiveness: High/Medium/Low

### Recommendation
Pursue / Monitor / Pass

### Next Steps
1. [Action item]

**Provide top 5-10 opportunities ranked by fit.**

**Before delivering, verify:**
- [ ] Eligibility requirements verified
- [ ] Deadlines are current
- [ ] Funder priorities confirmed
"""

        return await self.run_task(prompt)

    async def write_needs_statement(
        self,
        problem: str,
        target_population: str,
        data_sources: List[str] = None
    ) -> TaskResult:
        """Write a compelling needs statement"""
        await self.notify(f"Writing needs statement for: {problem[:50]}")

        prompt = f"""
Write needs statement for grant proposal:

**Problem:** {problem}
**Target Population:** {target_population}
{f"**Data sources:** {', '.join(data_sources)}" if data_sources else ""}

**Create needs statement following this structure:**

## Statement of Need

### The Problem
[Clear, compelling opening statement about the problem]

### Evidence
[3-5 data points from credible sources demonstrating the problem's scope]

- Statistic 1 (Source, Year)
- Statistic 2 (Source, Year)
- Statistic 3 (Source, Year)

### Impact
[Who is affected and how - make it personal and relatable]

### Current Gap
[What existing solutions are missing or inadequate]

### Urgency
[Why this must be addressed now - trends, deadlines, windows of opportunity]

### Our Unique Position
[Why we are positioned to address this need]

**Word count target:** 300-500 words
**Tone:** Urgent but hopeful, data-driven but human

**Before delivering, verify:**
- [ ] All statistics have citations
- [ ] Data is current (within 3-5 years)
- [ ] Aligns with funder priorities
"""

        return await self.run_task(prompt)

    async def create_budget(
        self,
        project_description: str,
        total_budget: float,
        duration_months: int = 12
    ) -> TaskResult:
        """Create grant budget with justification"""
        await self.notify(f"Creating budget: ${total_budget:,.0f}")

        prompt = f"""
Create grant budget and justification:

**Project:** {project_description}
**Total Budget:** ${total_budget:,.0f}
**Duration:** {duration_months} months

## Budget

### Personnel
| Position | FTE | Annual Salary | Effort | Cost |
|----------|-----|--------------|--------|------|
| Project Director | 1.0 | $X | X% | $X |
| ... | ... | ... | ... | ... |

**Subtotal Personnel:** $X

### Fringe Benefits (X% of salaries)
$X

### Travel
| Purpose | Cost |
|---------|------|
| ... | ... |

**Subtotal Travel:** $X

### Equipment (>$5,000/item)
| Item | Cost |
|------|------|
| ... | ... |

**Subtotal Equipment:** $X

### Supplies
| Category | Cost |
|----------|------|
| ... | ... |

**Subtotal Supplies:** $X

### Contractual
| Vendor/Purpose | Cost |
|----------------|------|
| ... | ... |

**Subtotal Contractual:** $X

### Other Direct Costs
| Item | Cost |
|------|------|
| ... | ... |

**Subtotal Other:** $X

### Indirect Costs (X% MTDC)
$X

## TOTAL: ${total_budget:,.0f}

## Budget Justification

### Personnel
[Detailed justification for each position]

### Fringe Benefits
[Rate and basis]

### Travel
[Purpose and breakdown of each trip]

### Equipment
[Why needed, alternatives considered]

### Supplies
[Categories and usage]

### Contractual
[Scope, selection process]

### Other
[Detailed breakdown]

### Indirect
[Rate agreement reference]

**Before delivering, verify:**
- [ ] Math adds up correctly
- [ ] No unallowable costs
- [ ] Justifications are specific
"""

        return await self.run_task(prompt)

    async def evaluate_opportunity(
        self,
        opportunity_name: str,
        rfp_url: str = None
    ) -> TaskResult:
        """Evaluate a specific funding opportunity"""
        prompt = f"""
Evaluate funding opportunity: {opportunity_name}
{f"RFP/FOA: {rfp_url}" if rfp_url else ""}

## Go/No-Go Analysis

### Basic Eligibility
| Requirement | Status | Notes |
|-------------|--------|-------|
| Organization type | ✓/✗ | ... |
| Geographic | ✓/✗ | ... |
| Prior relationship | ✓/✗ | ... |

### Strategic Fit
| Factor | Score (1-5) | Notes |
|--------|-------------|-------|
| Mission alignment | X | ... |
| Capacity to execute | X | ... |
| Competitive position | X | ... |
| Strategic value | X | ... |

**Total Score:** X/20

### Resource Requirements
- Proposal development: X hours
- Match/cost share: $X
- Personnel commitment: X FTE
- New capabilities needed: [List]

### Risks
| Risk | Impact | Likelihood | Mitigation |
|------|--------|------------|------------|
| ... | ... | ... | ... |

### Decision

**Recommendation:** Pursue / Pass / Request More Info

**Rationale:**
[Explanation]

**If pursuing, next steps:**
1. [Action with deadline]
2. [Action with deadline]
"""

        return await self.run_task(prompt)

    async def write_objectives(self, project_goals: str) -> TaskResult:
        """Write SMART objectives"""
        prompt = f"""
Write SMART objectives for:

{project_goals}

## Goals and Objectives

### Goal 1: [Broad goal statement]

**Objective 1.1:**
By [DATE], [TARGET NUMBER] [TARGET POPULATION] will [MEASURABLE OUTCOME]
as measured by [ASSESSMENT METHOD], resulting in [QUANTIFIABLE CHANGE].

**Objective 1.2:**
...

### Goal 2: [Broad goal statement]

**Objective 2.1:**
...

## Logic Model

| Inputs | Activities | Outputs | Short-term Outcomes | Long-term Outcomes |
|--------|------------|---------|---------------------|-------------------|
| ... | ... | ... | ... | ... |

## Evaluation Plan

| Objective | Indicator | Data Source | Frequency | Target |
|-----------|-----------|-------------|-----------|--------|
| 1.1 | ... | ... | ... | ... |

**Before delivering, verify:**
- [ ] All objectives are SMART
- [ ] Metrics are measurable
- [ ] Timeline is realistic
"""

        return await self.run_task(prompt)

    # ==================== DOCUMENT GENERATION METHODS ====================

    async def create_presentation(
        self,
        title: str,
        content: dict,
        template_path: str = None,
        output_path: str = None
    ) -> TaskResult:
        """Create a PowerPoint presentation from template or scratch"""
        await self.notify(f"Creating presentation: {title}")

        prompt = f"""
Create a PowerPoint presentation using python-pptx:

**Title:** {title}
**Template:** {template_path or "Create from scratch with professional styling"}
**Output:** {output_path or f"output/{title.lower().replace(' ', '-')}.pptx"}
**Content:** {content}

**Implementation:**
```python
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Load template or create new
{"prs = Presentation('" + template_path + "')" if template_path else "prs = Presentation()"}

# Set slide dimensions (16:9 widescreen)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add slides based on content structure
# For each section in content, create appropriate slide type:
# - Title slide: layout 0 or 6
# - Content slide: layout 1
# - Section header: layout 2
# - Two content: layout 3
# - Comparison: layout 4
# - Title only: layout 5
# - Blank: layout 6

# Apply consistent styling:
# - Font: Calibri or Arial
# - Title size: 44pt
# - Body size: 24pt
# - Colors: Use template colors or professional palette

# Save presentation
prs.save('{output_path or f"output/{title.lower().replace(' ', '-')}.pptx"}')
```

**Slide Structure Guidelines:**
1. Title slide with presentation title and subtitle
2. Agenda/overview slide
3. Content slides (one key message per slide)
4. Summary/conclusion slide
5. Q&A or contact slide

**Placeholder Replacement:**
If using template, replace placeholders:
- {{{{TITLE}}}} -> Presentation title
- {{{{DATE}}}} -> Current date
- {{{{AUTHOR}}}} -> Author name
- {{{{CONTENT}}}} -> Section content

**Before delivering, verify:**
- [ ] All slides render correctly
- [ ] Formatting is consistent
- [ ] Images are properly sized
- [ ] Text is readable
- [ ] File saves without errors
"""

        return await self.run_task(prompt)

    async def create_document(
        self,
        title: str,
        content: dict,
        template_path: str = None,
        output_path: str = None,
        doc_type: str = "report"
    ) -> TaskResult:
        """Create a Word document from template or scratch"""
        await self.notify(f"Creating document: {title}")

        prompt = f"""
Create a Word document using python-docx:

**Title:** {title}
**Type:** {doc_type}
**Template:** {template_path or "Create from scratch with professional styling"}
**Output:** {output_path or f"output/{title.lower().replace(' ', '-')}.docx"}
**Content:** {content}

**Implementation:**
```python
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE

# Load template or create new
{"doc = Document('" + template_path + "')" if template_path else "doc = Document()"}

# Set document properties
doc.core_properties.title = '{title}'
doc.core_properties.author = 'ARES Security'

# Add content based on document type
# Structure for {doc_type}:

{"# Report structure" if doc_type == "report" else ""}
{"# 1. Title page" if doc_type == "report" else ""}
{"# 2. Executive summary" if doc_type == "report" else ""}
{"# 3. Table of contents" if doc_type == "report" else ""}
{"# 4. Body sections with headings" if doc_type == "report" else ""}
{"# 5. Conclusions/recommendations" if doc_type == "report" else ""}
{"# 6. Appendices" if doc_type == "report" else ""}

{"# PRD structure" if doc_type == "prd" else ""}
{"# 1. Overview" if doc_type == "prd" else ""}
{"# 2. Goals and objectives" if doc_type == "prd" else ""}
{"# 3. User stories" if doc_type == "prd" else ""}
{"# 4. Requirements (functional/non-functional)" if doc_type == "prd" else ""}
{"# 5. Success metrics" if doc_type == "prd" else ""}

{"# Training plan structure" if doc_type == "training" else ""}
{"# 1. Training overview" if doc_type == "training" else ""}
{"# 2. Learning objectives" if doc_type == "training" else ""}
{"# 3. Curriculum outline" if doc_type == "training" else ""}
{"# 4. Schedule" if doc_type == "training" else ""}
{"# 5. Assessment criteria" if doc_type == "training" else ""}

# Apply styles
# - Heading 1: 16pt, Bold
# - Heading 2: 14pt, Bold
# - Normal: 11pt
# - Use consistent spacing

# Save document
doc.save('{output_path or f"output/{title.lower().replace(' ', '-')}.docx"}')
```

**Placeholder Replacement:**
If using template, replace placeholders in paragraphs and tables:
- {{{{TITLE}}}} -> Document title
- {{{{DATE}}}} -> Current date
- {{{{VERSION}}}} -> Version number
- {{{{AUTHOR}}}} -> Author name

**Before delivering, verify:**
- [ ] Document opens correctly
- [ ] Styles are consistent
- [ ] Headers/footers are set
- [ ] Page numbers work
- [ ] Table of contents updates
"""

        return await self.run_task(prompt)

    async def create_pdf(
        self,
        title: str,
        content: dict,
        output_path: str = None,
        from_docx: str = None
    ) -> TaskResult:
        """Create a PDF document"""
        await self.notify(f"Creating PDF: {title}")

        prompt = f"""
Create a PDF document:

**Title:** {title}
**Output:** {output_path or f"output/{title.lower().replace(' ', '-')}.pdf"}
**Content:** {content}
{"**Source DOCX:** " + from_docx if from_docx else ""}

**Implementation Options:**

{"**Option 1: Convert from DOCX (recommended if source exists)**" if from_docx else ""}
{'''```python
from docx2pdf import convert
convert("''' + from_docx + '''", "''' + (output_path or f"output/{title.lower().replace(' ', '-')}.pdf") + '''")
```''' if from_docx else ""}

**Option 2: Create with ReportLab:**
```python
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors

# Create document
doc = SimpleDocTemplate(
    '{output_path or f"output/{title.lower().replace(' ', '-')}.pdf"}',
    pagesize=letter,
    rightMargin=72,
    leftMargin=72,
    topMargin=72,
    bottomMargin=72
)

# Define styles
styles = getSampleStyleSheet()
title_style = ParagraphStyle(
    'CustomTitle',
    parent=styles['Heading1'],
    fontSize=24,
    spaceAfter=30,
    alignment=1  # Center
)
heading_style = ParagraphStyle(
    'CustomHeading',
    parent=styles['Heading2'],
    fontSize=14,
    spaceBefore=20,
    spaceAfter=10
)
body_style = styles['Normal']

# Build content
story = []
story.append(Paragraph('{title}', title_style))
story.append(Spacer(1, 0.5*inch))

# Add content sections
for section_title, section_content in content.items():
    story.append(Paragraph(section_title, heading_style))
    story.append(Paragraph(section_content, body_style))
    story.append(Spacer(1, 0.25*inch))

# Build PDF
doc.build(story)
```

**Option 3: Create with PyPDF2 (for merging/modifying existing PDFs):**
```python
from PyPDF2 import PdfMerger, PdfReader, PdfWriter

# Merge multiple PDFs
merger = PdfMerger()
merger.append('file1.pdf')
merger.append('file2.pdf')
merger.write('{output_path or f"output/{title.lower().replace(' ', '-')}.pdf"}')
merger.close()
```

**Before delivering, verify:**
- [ ] PDF opens correctly
- [ ] All pages render
- [ ] Text is selectable
- [ ] Images are clear
- [ ] File size is reasonable
"""

        return await self.run_task(prompt)

    async def apply_template(
        self,
        template_path: str,
        data: dict,
        output_path: str
    ) -> TaskResult:
        """Apply data to a document template (PPTX, DOCX, or PDF)"""
        await self.notify(f"Applying template: {template_path}")

        # Determine file type
        file_ext = template_path.split('.')[-1].lower()

        prompt = f"""
Apply data to template:

**Template:** {template_path}
**Output:** {output_path}
**Data:** {data}
**Format:** {file_ext.upper()}

**Placeholder Pattern:**
Templates use {{{{PLACEHOLDER}}}} syntax for variable replacement.

**Implementation for {file_ext.upper()}:**

{'''```python
from pptx import Presentation

prs = Presentation("''' + template_path + '''")

data = ''' + str(data) + '''

# Iterate through slides and shapes
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in data.items():
                        placeholder = "{{" + key + "}}"
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(value))

prs.save("''' + output_path + '''")
```''' if file_ext == 'pptx' else ''}

{'''```python
from docx import Document

doc = Document("''' + template_path + '''")

data = ''' + str(data) + '''

# Replace in paragraphs
for paragraph in doc.paragraphs:
    for key, value in data.items():
        placeholder = "{{" + key + "}}"
        if placeholder in paragraph.text:
            for run in paragraph.runs:
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, str(value))

# Replace in tables
for table in doc.tables:
    for row in table.rows:
        for cell in row.cells:
            for key, value in data.items():
                placeholder = "{{" + key + "}}"
                if placeholder in cell.text:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(value))

doc.save("''' + output_path + '''")
```''' if file_ext == 'docx' else ''}

**Before delivering, verify:**
- [ ] All placeholders replaced
- [ ] Formatting preserved
- [ ] No {{{{PLACEHOLDER}}}} text remains
- [ ] Output file opens correctly
"""

        return await self.run_task(prompt)

    async def create_training_plan(
        self,
        title: str,
        objectives: List[str],
        duration: str,
        audience: str,
        output_format: str = "pptx"
    ) -> TaskResult:
        """Create a comprehensive training plan document"""
        await self.notify(f"Creating training plan: {title}")

        prompt = f"""
Create a training plan:

**Title:** {title}
**Learning Objectives:** {objectives}
**Duration:** {duration}
**Target Audience:** {audience}
**Output Format:** {output_format.upper()}

**Training Plan Structure:**

## 1. Training Overview
- Purpose and scope
- Target audience description
- Prerequisites

## 2. Learning Objectives
By the end of this training, participants will be able to:
{chr(10).join(f"- {obj}" for obj in objectives)}

## 3. Curriculum Outline

| Module | Topic | Duration | Method |
|--------|-------|----------|--------|
| 1 | Introduction & Overview | X hrs | Lecture |
| 2 | Core Concepts | X hrs | Interactive |
| 3 | Hands-on Practice | X hrs | Lab |
| 4 | Advanced Topics | X hrs | Workshop |
| 5 | Assessment & Review | X hrs | Exam/Project |

## 4. Detailed Module Breakdown

### Module 1: [Title]
- Topics covered
- Activities
- Materials needed
- Assessment criteria

[Repeat for each module]

## 5. Training Schedule

| Day | Time | Module | Instructor |
|-----|------|--------|------------|
| 1 | 09:00-12:00 | Module 1 | TBD |
| 1 | 13:00-17:00 | Module 2 | TBD |
[Continue...]

## 6. Assessment & Evaluation
- Knowledge checks
- Practical assessments
- Certification criteria (if applicable)

## 7. Resources & Materials
- Training materials list
- Equipment needed
- Reference documents

**Output as {output_format.upper()}:**
{"Create PowerPoint with one module per section, visual aids, and speaker notes" if output_format == "pptx" else ""}
{"Create Word document with professional formatting, headers, and table of contents" if output_format == "docx" else ""}
{"Create PDF with clean layout and print-ready formatting" if output_format == "pdf" else ""}

Save to: output/{title.lower().replace(' ', '-')}-training-plan.{output_format}

**Before delivering, verify:**
- [ ] All modules are complete
- [ ] Timeline is realistic
- [ ] Objectives map to assessments
- [ ] Materials list is comprehensive
"""

        return await self.run_task(prompt)

    async def create_prd(
        self,
        product_name: str,
        problem_statement: str,
        target_users: str,
        output_format: str = "docx"
    ) -> TaskResult:
        """Create a Product Requirements Document"""
        await self.notify(f"Creating PRD: {product_name}")

        prompt = f"""
Create a Product Requirements Document (PRD):

**Product:** {product_name}
**Problem:** {problem_statement}
**Target Users:** {target_users}
**Output Format:** {output_format.upper()}

**PRD Structure:**

# Product Requirements Document: {product_name}

## 1. Overview
### 1.1 Purpose
### 1.2 Scope
### 1.3 Definitions & Acronyms

## 2. Problem Statement
{problem_statement}

### 2.1 Current State
### 2.2 Pain Points
### 2.3 Opportunity

## 3. Goals & Objectives
### 3.1 Business Goals
### 3.2 User Goals
### 3.3 Success Metrics (KPIs)

## 4. Target Users
{target_users}

### 4.1 User Personas
### 4.2 User Journey Maps

## 5. Requirements

### 5.1 Functional Requirements
| ID | Requirement | Priority | Status |
|----|-------------|----------|--------|
| FR-001 | [Requirement] | Must Have | Draft |

### 5.2 Non-Functional Requirements
| ID | Requirement | Category | Priority |
|----|-------------|----------|----------|
| NFR-001 | [Requirement] | Performance | Must Have |

## 6. User Stories
### Epic 1: [Name]
- As a [user], I want to [action] so that [benefit]
  - Acceptance Criteria:
    - [ ] Criterion 1
    - [ ] Criterion 2

## 7. Design & UX
### 7.1 Wireframes (references)
### 7.2 User Flows

## 8. Technical Considerations
### 8.1 Architecture Overview
### 8.2 Integration Points
### 8.3 Security Requirements

## 9. Timeline & Milestones
| Phase | Milestone | Target Date |
|-------|-----------|-------------|
| Phase 1 | MVP | TBD |

## 10. Risks & Mitigations
| Risk | Impact | Likelihood | Mitigation |
|------|--------|------------|------------|

## 11. Appendix
### A. References
### B. Change Log

Save to: output/{product_name.lower().replace(' ', '-')}-prd.{output_format}

**Before delivering, verify:**
- [ ] All sections complete
- [ ] Requirements are measurable
- [ ] User stories have acceptance criteria
- [ ] Priorities are assigned
"""

        return await self.run_task(prompt)

    # ==================== COMBINED METHODS ====================

    async def create_documentation_suite(
        self,
        project_name: str,
        include_api: bool = True,
        include_architecture: bool = True,
        include_user_guide: bool = True
    ) -> TaskResult:
        """Create a complete documentation suite for a project"""
        await self.notify(f"Creating documentation suite for: {project_name}")

        prompt = f"""
Create complete documentation suite for: {project_name}

**Include:**
- API Documentation: {include_api}
- Architecture Documentation: {include_architecture}
- User Guide: {include_user_guide}

## 1. README.md
Create comprehensive README with:
- Project overview
- Quick start guide
- Installation instructions
- Basic usage examples
- Links to detailed docs

{"## 2. API Documentation" if include_api else ""}
{'''
Create docs/api/README.md with:
- Authentication guide
- Endpoint reference
- Request/response examples
- Error handling guide
''' if include_api else ""}

{"## 3. Architecture Documentation" if include_architecture else ""}
{'''
Create docs/architecture/README.md with:
- System overview diagram
- Component descriptions
- Data flow documentation
- Security considerations
''' if include_architecture else ""}

{"## 4. User Guide" if include_user_guide else ""}
{'''
Create docs/guide/README.md with:
- Getting started tutorial
- Common use cases
- Configuration reference
- Troubleshooting guide
''' if include_user_guide else ""}

## 5. Documentation Index
Create docs/README.md linking all documentation

**Before delivering, verify:**
- [ ] All docs are consistent in style
- [ ] Cross-references are valid
- [ ] Examples are tested
- [ ] Navigation is clear
"""

        return await self.run_task(prompt)

    async def work(self, task: str) -> TaskResult:
        """General documentation/writing work"""
        await self.notify(f"Starting: {task[:50]}...")
        return await self.run_task(task)


async def main():
    """CLI entry point"""
    import argparse
    import json

    parser = argparse.ArgumentParser(description="Valentina - Technical Writer & Content Strategist")

    # Documentation arguments
    parser.add_argument("--feature", type=str, help="Document feature")
    parser.add_argument("--api", type=str, help="Document API endpoint")
    parser.add_argument("--diagram", type=str, help="Create diagram for subject")
    parser.add_argument("--type", type=str, default="flow", help="Diagram type")
    parser.add_argument("--readme", type=str, nargs="?", const=".", help="Update README")
    parser.add_argument("--architecture", type=str, help="Create architecture doc")
    parser.add_argument("--suite", type=str, help="Create documentation suite for project")

    # Grant writing arguments
    parser.add_argument("--research", type=str, help="Research funding for project type")
    parser.add_argument("--budget-range", type=str, help="Budget range for research")
    parser.add_argument("--needs", type=str, help="Write needs statement for problem")
    parser.add_argument("--population", type=str, help="Target population for needs")
    parser.add_argument("--create-budget", type=str, help="Create budget for project")
    parser.add_argument("--amount", type=float, help="Budget amount")
    parser.add_argument("--duration", type=int, default=12, help="Duration in months")
    parser.add_argument("--evaluate", type=str, help="Evaluate opportunity")
    parser.add_argument("--objectives", type=str, help="Write objectives for goals")

    # Document generation arguments
    parser.add_argument("--presentation", type=str, help="Create PowerPoint presentation")
    parser.add_argument("--document", type=str, help="Create Word document")
    parser.add_argument("--doc-type", type=str, default="report", help="Document type: report, prd, training")
    parser.add_argument("--pdf", type=str, help="Create PDF document")
    parser.add_argument("--from-docx", type=str, help="Convert DOCX to PDF")
    parser.add_argument("--template", type=str, help="Template file path")
    parser.add_argument("--output", type=str, help="Output file path")
    parser.add_argument("--data", type=str, help="JSON data for template placeholders")
    parser.add_argument("--apply-template", type=str, help="Apply data to template")
    parser.add_argument("--training-plan", type=str, help="Create training plan")
    parser.add_argument("--prd", type=str, help="Create PRD for product")
    parser.add_argument("--problem", type=str, help="Problem statement for PRD")
    parser.add_argument("--users", type=str, help="Target users for PRD")
    parser.add_argument("--format", type=str, default="pptx", help="Output format: pptx, docx, pdf")

    # General
    parser.add_argument("--task", type=str, help="Run general task")
    parser.add_argument("--status", action="store_true", help="Show status")

    args = parser.parse_args()

    agent = ValentinaAgent()

    if args.status:
        print(json.dumps(agent.get_status(), indent=2))
        return

    # Documentation operations
    if args.feature:
        result = await agent.document_feature(args.feature)
        print(result.output)
        return

    if args.api:
        result = await agent.document_api(args.api)
        print(result.output)
        return

    if args.diagram:
        result = await agent.create_diagram(args.diagram, args.type)
        print(result.output)
        return

    if args.readme is not None:
        result = await agent.update_readme(args.readme)
        print(result.output)
        return

    if args.architecture:
        result = await agent.create_architecture_doc(args.architecture)
        print(result.output)
        return

    if args.suite:
        result = await agent.create_documentation_suite(args.suite)
        print(result.output)
        return

    # Grant writing operations
    if args.research:
        result = await agent.research_funding(args.research, args.budget_range)
        print(result.output)
        return

    if args.needs and args.population:
        result = await agent.write_needs_statement(args.needs, args.population)
        print(result.output)
        return

    if args.create_budget and args.amount:
        result = await agent.create_budget(args.create_budget, args.amount, args.duration)
        print(result.output)
        return

    if args.evaluate:
        result = await agent.evaluate_opportunity(args.evaluate)
        print(result.output)
        return

    if args.objectives:
        result = await agent.write_objectives(args.objectives)
        print(result.output)
        return

    # Document generation operations
    if args.presentation:
        content = json.loads(args.data) if args.data else {}
        result = await agent.create_presentation(
            args.presentation, content, args.template, args.output
        )
        print(result.output)
        return

    if args.document:
        content = json.loads(args.data) if args.data else {}
        result = await agent.create_document(
            args.document, content, args.template, args.output, args.doc_type
        )
        print(result.output)
        return

    if args.pdf:
        content = json.loads(args.data) if args.data else {}
        result = await agent.create_pdf(args.pdf, content, args.output, args.from_docx)
        print(result.output)
        return

    if args.apply_template and args.data:
        data = json.loads(args.data)
        result = await agent.apply_template(args.apply_template, data, args.output)
        print(result.output)
        return

    if args.training_plan:
        objectives = args.objectives.split(',') if args.objectives else []
        result = await agent.create_training_plan(
            args.training_plan, objectives, args.duration, args.population or "General", args.format
        )
        print(result.output)
        return

    if args.prd:
        result = await agent.create_prd(
            args.prd, args.problem or "", args.users or "", args.format
        )
        print(result.output)
        return

    if args.task:
        result = await agent.work(args.task)
        print(result.output)
        return

    print("Valentina - Technical Writer & Content Strategist")
    print("==================================================")
    print("Use --help for options")


if __name__ == "__main__":
    asyncio.run(main())
